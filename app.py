# Import required libraries
import streamlit as st
import openai
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO

# Set up OpenAI API Key
openai.api_key = "your  key"

# Define available themes and fonts
themes = {
    "Classic": RGBColor(0, 0, 0),         # Black text
    "Modern": RGBColor(255, 0, 0),        # Red text
    "Professional": RGBColor(0, 0, 255)   # Blue text
}
font_families = ["Calibri", "Arial", "Times New Roman", "Verdana"]

# Function to generate content for structured slides
def generate_slide_content(topic, slide_count, bullet_count, model="gpt-3.5-turbo"):
    slides_content = []

    # Generate introduction slide content
    intro_prompt = f"Write a brief introduction for a presentation on {topic}."
    intro_response = openai.ChatCompletion.create(
        model=model,
        messages=[{"role": "user", "content": intro_prompt}],
        max_tokens=1000
    )
    introduction = intro_response['choices'][0]['message']['content'].strip()
    slides_content.append([introduction])  # First slide is the introduction

    # Generate index slide content (subtopics)
    index_prompt = f"List the main sections or subtopics for a presentation on {topic}."
    index_response = openai.ChatCompletion.create(
        model=model,
        messages=[{"role": "user", "content": index_prompt}],
        max_tokens=100
    )
    index_content = index_response['choices'][0]['message']['content'].strip().split('\n')
    slides_content.append(index_content)  # Second slide is the index

    # Generate content for each subtopic slide
    for i in range(slide_count):
        prompt = f"Generate up to {bullet_count} bullet points for slide {i+1} of a presentation about {topic}."
        response = openai.ChatCompletion.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=150
        )
        content = response['choices'][0]['message']['content'].strip()
        
        # Ensure only the specified number of bullet points per slide
        bullet_points = content.split('\n')
        slides_content.append(bullet_points[:bullet_count])  # Add the content for each subtopic slide

    return slides_content

# Function to create a PowerPoint presentation with theme and font customization
def create_ppt(topic, slides_content, theme_color, font_family, font_size, use_bullets):
    ppt = Presentation()

    for idx, content in enumerate(slides_content):
        slide_layout = ppt.slide_layouts[5]  # Blank layout for custom positioning
        slide = ppt.slides.add_slide(slide_layout)

        # Apply gradient-like background with light and darker blue rectangles
        background_shape1 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        background_shape1.fill.solid()
        background_shape1.fill.fore_color.rgb = RGBColor(220, 220, 255)  # Light blue
        
        background_shape2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.2), Inches(9.6), Inches(7.1)
        )
        background_shape2.fill.solid()
        background_shape2.fill.fore_color.rgb = RGBColor(200, 200, 240)  # Slightly darker blue

        # Add border around the slide
        border_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.2), Inches(9.6), Inches(7.1)
        )
        border_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        border_shape.line.width = Pt(2)

        # Set slide title to be centered at the top with underline
        slide_title = topic if idx >= 2 else ("Introduction" if idx == 0 else "Index of Topics")
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.bold = True
        title_paragraph.font.size = Pt(32)  # Larger font for H1-style heading
        title_paragraph.font.color.rgb = theme_color
        title_paragraph.font.name = font_family
        title_paragraph.font.underline = True  # Underline for the heading
        title_paragraph.alignment = PP_ALIGN.CENTER  # Center alignment for heading

        # Add bullet points/content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        content_box.text_frame.word_wrap = True  # Enable word wrap
        for point in content:
            p = content_box.text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(font_size)
            p.font.color.rgb = theme_color
            p.font.name = font_family
            p.level = 0 if use_bullets else None  # Set as bullet if enabled

    return ppt

# Streamlit UI
st.title("Customizable AI-Powered PowerPoint Generator")

# User input for topic and slide count
topic = st.text_input("Enter the topic of the presentation:")
slide_count = st.number_input("Number of subtopic slides", min_value=1, max_value=20, value=5)

# Model selection
model_choice = st.selectbox("Select AI model:", ["gpt-3.5-turbo", "gpt-4"])

# Theme, font, bullet options, and bullet count
theme_choice = st.selectbox("Select theme color:", list(themes.keys()))
font_choice = st.selectbox("Select font family:", font_families)
font_size = st.slider("Select font size for slide text:", 10, 32, 18)  # Font size slider
use_bullets = st.checkbox("Use bullet points", value=True)  # Option to enable/disable bullets
bullet_count = st.number_input("Max number of bullet points per slide", min_value=1, max_value=10, value=5)

if st.button("Generate Presentation"):
    if topic:
        with st.spinner("Generating slides..."):
            slides_content = generate_slide_content(topic, slide_count, bullet_count, model=model_choice)
            ppt = create_ppt(topic, slides_content, themes[theme_choice], font_choice, font_size, use_bullets)
            
            # Save the presentation to a BytesIO stream instead of a file path
            ppt_bytes = BytesIO()
            ppt.save(ppt_bytes)
            ppt_bytes.seek(0)  # Rewind to the start of the BytesIO buffer

            # Provide download link
            st.download_button(
                label="Download Presentation",
                data=ppt_bytes,
                file_name="presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    else:
        st.error("Please enter a topic for the presentation.")
